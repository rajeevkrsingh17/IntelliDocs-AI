from pptx import Presentation


def extract_pptx_text(pptx_path):
    """
    Extract text from a PPTX presentation and return text with metadata.
    """

    presentation = Presentation(pptx_path)

    slides_text = []
    page_count = len(presentation.slides)

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():
                    slides_text.append(shape.text)

    full_text = "\n".join(slides_text)
    return {"text": full_text, "pages": page_count}